import pandas as pd
import numpy as np
from PyQt5.QtGui import QFont, QIcon, QPainter, QPixmap, QColor
from PyQt5.QtCore import Qt, QMargins, QTranslator, QLocale, QLibraryInfo
from PyQt5.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, QSizePolicy,
                             QGridLayout, QLabel, QColorDialog, QComboBox, QCheckBox,
                             QPushButton, QFileDialog, QMessageBox, QLineEdit, QGraphicsView, QTableWidget, QTableWidgetItem)
from PyQt5 import uic, QtWidgets, QtCore, QtWebEngineWidgets
import sys
import datetime
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import plotly.express as px
import plotly

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

form_class = uic.loadUiType("LGSTA_V2.0_test6.ui")[0]

class WindowClass(QMainWindow, form_class) :
    def __init__(self, *args, **kwargs) :
        super(WindowClass, self).__init__(*args, **kwargs)
        self.setupUi(self)
        self.browser = QtWebEngineWidgets.QWebEngineView()
        self.browser1 = QtWebEngineWidgets.QWebEngineView()
        self.browser2 = QtWebEngineWidgets.QWebEngineView()
        self.browser3 = QtWebEngineWidgets.QWebEngineView()
        self.browser4 = QtWebEngineWidgets.QWebEngineView()

        # 초기값 ----------------------------------------------------
        self.comp = pd.DataFrame()
        self.comp_add = pd.DataFrame()
        self.mv0_odu = pd.DataFrame()
        self.dfs = pd.DataFrame()
        self.r_colist_rev = []
        self.s_colist_rev = []
        self.Hz_col_rev = []
        self.lim_C_rev = []
        self.lim_STS_rev = []
        self.tdiff = 0
        self.model = "-"
        self.condition = "-"
        self.location = "-"
        self.sample = "-"
        self.day =""
        self.year =""
        self.month = ""
        self.place =""
        self.max_df1 = pd.DataFrame()
        #------------------------------------------------------------

        self.line_model.textChanged.connect(self.getTextFunc1)
        self.line_place.textChanged.connect(self.getTextFunc4)
        self.line3.textChanged.connect(self.getTextFunc3)

        self.bt1.clicked.connect(self.load_stress)
        self.bt2.clicked.connect(self.load_mv)

        self.cb1.setChecked(False)
        self.cb1.stateChanged.connect(lambda: self.btnstate(self.cb1))
        self.cb2.setChecked(False)
        self.cb2.stateChanged.connect(lambda: self.btnstate(self.cb2))
        self.cb3.setChecked(False)
        self.cb3.stateChanged.connect(lambda: self.btnstate(self.cb3))
        self.cb4.setChecked(False)
        self.cb4.stateChanged.connect(lambda: self.btnstate(self.cb4))
        self.cb5.setChecked(False)
        self.cb5.stateChanged.connect(lambda: self.btnstate(self.cb5))

        self.table1.setColumnWidth(0, 80)

        for i in range(6):
            self.table_max.setColumnWidth(i, 80)

        self.bt4.clicked.connect(self.show_graph)
        self.bt4.clicked.connect(self.max_list_show)
        self.bt5.clicked.connect(self.capture)
        self.bt6.clicked.connect(self.makeppt)
        self.bt7.clicked.connect(self.disp_imglist)
        self.bt8.clicked.connect(self.opendir)
        self.bt9.clicked.connect(self.opendir_rep)

        self.comboBox.activated.connect(self.comboGetText)
        self.comboBox_2.activated.connect(self.combo2GetText)
        self.comboBox_3.activated.connect(self.combo3GetText)
        self.comboBox_4.activated.connect(self.combo4GetText)
        self.comboBox_5.activated.connect(self.combo5GetText)
        self.comboBox_6.activated.connect(self.combo6GetText)

        self.graph_b0.addWidget(self.browser)
        self.graph_b1.addWidget(self.browser1)
        self.graph_b2.addWidget(self.browser2)
        self.graph_b3.addWidget(self.browser3)
        self.graph_b4.addWidget(self.browser4)

    def comboGetText(self):
        self.group = str(self.comboBox.currentText())
    def combo2GetText(self):
        self.condition = str(self.comboBox_2.currentText())
        if self.condition == "--선택--":
            self.condition = "--"
    def combo3GetText(self):
        self.sample = str(self.comboBox_3.currentText())
    def combo4GetText(self):
        self.year = self.comboBox_4.currentText()
    def combo5GetText(self):
        self.month = self.comboBox_5.currentText()
    def combo6GetText(self):
        self.day = self.comboBox_6.currentText()

    def opendir(self):
        path="./CapturedGraphs"
        if not os.path.isdir(path):
            os.mkdir(path)
        path=os.path.realpath(path)
        os.startfile(path)

    def opendir_rep(self):
        path = "./report"
        if not os.path.isdir(path):
            os.mkdir(path)
        path = os.path.realpath(path)
        os.startfile(path)

# table 체크박스 filter coding 시작------------------------------------------------------------------------------------------------------------------------
    def disp_imglist(self):
        path_dir = "./CapturedGraphs/"
        files = os.listdir(path_dir)
        for i in range(0, len(files)):
            for j in range(0, len(files)):
                if datetime.datetime.fromtimestamp(os.path.getmtime(path_dir+files[i])) < datetime.datetime.fromtimestamp(os.path.getmtime(path_dir+files[j])):
                    (files[i], files[j]) = (files[j], files[i])
        i=0
        self.table1.setRowCount(len(files))
        self.table1.setColumnCount(1)
        self.table1.setColumnWidth(0, 200)
        self.table1.setHorizontalHeaderLabels(["Condition"])
        self.table1.horizontalHeaderItem(0).setTextAlignment(Qt.AlignHCenter)

        self.file_list = files

        for list in files:
            self.table1.setItem(i, 0, QTableWidgetItem(list))
            i = i + 1

        self.tableHeader = self.table1.horizontalHeader()
        self.tableHeader.sectionClicked.connect(self.columnfilterclicked)
        self.keywords = dict([(i, []) for i in range(self.table1.columnCount())])
        self.checkBoxs = []
        self.col = None

    def slotSelect(self, state):
        for checkbox in self.checkBoxs:
            checkbox.setChecked(QtCore.Qt.Checked == state)

    def menuClose(self):
        self.keywords[self.col] = []
        for element in self.checkBoxs:
            if element.isChecked():
                self.keywords[self.col].append(element.text())
        self.filterdata()
        self.menu.close()

    def clearFilter(self):
        if self.table1.rowCount() > 0:
            for i in range(self.table1.rowCount()):
                self.table1.setRowHidden(i, False)

    def filterdata(self):
        columnsShow = dict([(i, True) for i in range(self.table1.rowCount())])

        for i in range(self.table1.rowCount()):
            for j in range(self.table1.columnCount()):
                item = self.table1.item(i, j)
                if self.keywords[j]:
                    if item.text() not in self.keywords[j]:
                        columnsShow[i] = False
        for key in columnsShow:
            self.table1.setRowHidden(key, not columnsShow[key])

    def columnfilterclicked(self, index):
        self.menu = QtWidgets.QMenu(self)
        self.col = index

        data_unique = []
        self.checkBoxs = []

        checkBox = QtWidgets.QCheckBox("Select all", self.menu)
        checkableAction = QtWidgets.QWidgetAction(self.menu)
        checkableAction.setDefaultWidget(checkBox)
        self.menu.addAction(checkableAction)
        checkBox.setChecked(True)
        checkBox.stateChanged.connect(self.slotSelect)

        for i in range(self.table1.rowCount()):
            if not self.table1.isRowHidden(i):
                item = self.table1.item(i, index)
                if item.text() not in data_unique:
                    data_unique.append(item.text())
                    checkBox = QtWidgets.QCheckBox(item.text(), self.menu)
                    checkBox.setChecked(True)
                    checkableAction = QtWidgets.QWidgetAction(self.menu)
                    checkableAction.setDefaultWidget(checkBox)
                    self.menu.addAction(checkableAction)
                    self.checkBoxs.append(checkBox)

        btn = QtWidgets.QDialogButtonBox(QtWidgets.QDialogButtonBox.Ok | QtWidgets.QDialogButtonBox.Cancel,
                                             QtCore.Qt.Horizontal, self.menu)
        btn.accepted.connect(self.menuClose)
        btn.rejected.connect(self.menu.close)
        checkableAction = QtWidgets.QWidgetAction(self.menu)
        checkableAction.setDefaultWidget(btn)
        self.menu.addAction(checkableAction)

        headerPos = self.table1.mapToGlobal(self.tableHeader.pos())

        posY = headerPos.y() + self.tableHeader.height()
        posX = headerPos.x() + self.tableHeader.sectionPosition(index)
        self.menu.exec_(QtCore.QPoint(posX, posY))

# table 체크박스 filter coding 끝------------------------------------------------------------------------------------------------------------------------

    def getTextFunc1(self) :
        self.model = str(self.line_model.text())

    def getTextFunc3(self):
        try:
            self.tdiff = int(self.line3.text())
        except ValueError:
            pass
    def getTextFunc4(self):
        self.place = str(self.line_place.text())

    def load_stress(self):
        try:
            fname = QFileDialog.getOpenFileName(self, 'Open file', './data')
            if fname[0]:
                dfs = pd.read_csv(fname[0], encoding='cp949')
                dfs = dfs.drop(dfs.index[0])
                dfs_col = int((dfs.shape[1] + 1) / 2)
                dfs = dfs.rename({'CH Name': 'sec'}, axis=1)

                S_col = dfs.columns

                s_colist = []
                r_colist = []

                for i in range(1, dfs_col):
                    s_colist.append(S_col[i])


                for j in range(dfs_col, dfs.shape[1]):
                    r_colist.append(S_col[j])

                number_place = ((len(dfs.columns) - 1) / 2)
                number_place = int(number_place)
                pipe_place = dfs.columns[1:(number_place + 1)]

                dfs['sec'] = (dfs['sec'].astype(int)-1)
                dfs[pipe_place] = dfs[pipe_place].astype(float)

                ####################################################
                odd = dfs[1::2]
                odd = odd.reset_index(drop=True)
                even = dfs[0::2]
                even = even.reset_index(drop=True)

                if len(odd.index) == len(even.index):
                    pass
                else:
                    even = even[:(len(even.index)-1)]

                dfs_reduce = pd.DataFrame()

                for i in S_col:
                    dfs_reduce[i] = odd[i].where(odd[i] > even[i],even[i])
                dfs_reduce['sec'] = even['sec']

                dfs_reduce['Spec_C'] = 15
                dfs_reduce['Spec_STS'] =90
                ######################################################
                self.dfs = dfs_reduce
                self.pipe_place = pipe_place
                self.s_colist = s_colist
                self.r_colist = r_colist
                self.lim_C = ['Spec_C']
                self.lim_STS = ['Spec_STS']

                self.progressBar.setValue(100)

                QMessageBox.about(self, 'OK', '업로드가 완료되었습니다.')

            else:
                QMessageBox.about(self, 'Warning', '파일을 선택하지 않았습니다.')

        except UnicodeDecodeError:
            QMessageBox.about(self, 'Warning', '응력변환 원본 파일이 아닙니다. 원본 파일을 선택하세요. [사용 방법]을 참조하세요.')

    def load_mv(self):
        try:
            fname = QFileDialog.getOpenFileName(self, 'Open file', './data')
            if fname[0]:
                # 파일 읽기
                mv_odu = pd.read_csv(fname[0], skiprows=6)
                mv_odu['mv_sec'] = (mv_odu['siter'] - 1) * 2
                comp = mv_odu[['mv_sec', 'INV 현재']]

                float32 = ['Error', 'ErrorUnit', '압축비', '설정온도차', 'INV 목표', 'INV 현재', 'FAN 목표', 'FAN 현재', 'FAN2 목표',
                           'FAN2 현재', '공기온도', '흡입온도', '응축온도', '증발온도', 'INV 토출온도', '열교환기온도', '열교환기출구', 'Comp Heatsink',
                           'INV토출과열도', 'IHEX 입구', 'IHEX 출구']
                int16 = ['siter', '운전모드', '목표고압', '목표저압', 'INV1토출목표', '현재고압', '현재저압', 'COMP기준', 'PC CTRL', '운전상태 유지',
                         '4WAY', 'HOT GAS', 'SUMP HEAT', 'Cooling Fan', '오일분리기', 'POWERRELAY', 'PFC', 'PFC CTRL', 'COP',
                         'DEFROSTING', 'OIL REFUREL', 'OIL RETURN', 'FAN1 H', 'FAN1 L', 'FAN2 H', 'FAN2 L', 'Main EEV',
                         'COMP운전', 'Sub EEV', 'VI EEV1', 'mv_sec']
                mv_odu_list = float32 + int16

                mv_odu.loc[:, float32] = mv_odu.loc[:, float32].astype('float32')
                mv_odu.loc[:, int16] = mv_odu.loc[:, int16].astype('int16')

                mv0_odu = mv_odu[mv_odu_list]
                mv0_odu = pd.merge(mv0_odu, mv_odu['mv_sec'])

                self.mv0_odu = mv0_odu

                # self.comp_add.loc[:, 'mv_sec'] = comp['mv_sec'] + 1
                # for i in range(len(comp) - 1):
                #     self.comp_add.loc[i, 'INV 현재'] = round((comp.loc[i, 'INV 현재'] + comp.loc[i + 1, 'INV 현재']) / 2)
                # self.comp_dt1 = pd.concat((comp, self.comp_add), axis=0).sort_values(by=['mv_sec']).dropna(axis=0)
                # print(self.comp_dt1)
                self.Hz_col = ['INV 현재']
                self.comp = comp

                # ODU1 파일명 확인하기
                filename = os.path.splitext(fname[0])[0]
                mv_name = filename[:-5]
                # 전체 mv데이터 불러오기
                idu1 = mv_name + "_IDU1.csv"
                odu2 = mv_name + "_ODU2.csv"
                emctl = mv_name + "_EMCTL.csv"
                # 파일 존재여부 확인
                if os.path.isfile(idu1):
                    data_IDU1 = pd.read_csv(idu1, skiprows=6)
                else:
                    QMessageBox.about(self, 'Warning', 'IDU1 파일이 없습니다.')
                if os.path.isfile(odu2):
                    data_ODU2 = pd.read_csv(odu2, skiprows=6)
                else:
                    QMessageBox.about(self, 'Warning', 'ODU2 파일이 없습니다.')
                if os.path.isfile(emctl):
                    data_EMCTL = pd.read_csv(emctl, skiprows=6)
                else:
                    QMessageBox.about(self, 'Warning', 'EMCTL 파일이 없습니다.')

                self.progressBar_2.setValue(100)
                QMessageBox.about(self, 'OK', '업로드가 완료되었습니다.')
            else:
                QMessageBox.about(self, 'Warning', '파일을 선택하지 않았습니다.')
        except (KeyError, UnicodeDecodeError, NameError, pd.errors.EmptyDataError, pd.errors.ParserError):
            QMessageBox.about(self, 'Warning', 'ODU1 MV 파일이 아닙니다. ODU1 MV 원본파일을 선택하십시오.')



    def btnstate(self, cb):
        if cb.text() == "응력의 Comp Hz":
            if cb.isChecked() == True:
                if self.dfs.empty:
                    QMessageBox.about(self, 'Warning', '응력 데이터가 없습니다. 응력 데이터를 입력하세요.')
                else:
                    self.r_colist_rev = self.r_colist
            else:
                self.r_colist_rev = []
        if cb.text() == "MV의 Comp Hz":
            if cb.isChecked() == True:
                if self.comp.empty:
                    QMessageBox.about(self, 'Warning', 'MV 데이터가 없습니다. MV 데이터를 입력하세요.')
                else:
                    self.Hz_col_rev = self.Hz_col
            else:
                self.Hz_col_rev = []
        if cb.text() == "CU (15MPa)":
            if cb.isChecked() == True:
                if self.dfs.empty:
                    QMessageBox.about(self, 'Warning', '응력 데이터가 없습니다. 응력 데이터를 입력하세요.')
                else:
                    self.lim_C_rev = self.lim_C
            else:
                self.lim_C_rev = []
        if cb.text() == "STS (90MPa)":
            if cb.isChecked() == True:
                if self.dfs.empty:
                    QMessageBox.about(self, 'Warning', '응력 데이터가 없습니다. 응력 데이터를 입력하세요.')
                else:
                    self.lim_STS_rev = self.lim_STS
            else:
                self.lim_STS_rev = []
        if cb.text() == "응력 데이터":
            if cb.isChecked() == True:
                if self.dfs.empty:
                    QMessageBox.about(self, 'Warning', '응력 데이터가 없습니다. 응력 데이터를 입력하세요.')
                else:
                    self.s_colist_rev = self.s_colist
            else:
                self.s_colist_rev = []

    def show_graph(self):
        if self.dfs.empty:
            QMessageBox.about(self, 'Warning', '응력 데이터가 없습니다. 응력 데이터를 입력하세요.')
        else:

            if self.mv0_odu.empty and self.tdiff == 0 :
                # s_colist : 응력값
                # r_colist : 응력데이터의 comp rps값

                fig0 = make_subplots(specs=[[{"secondary_y": True}]])
                fig0.update_layout(margin=dict(t=40, b=40, l=10, r=10), template="plotly_white",)

                for i in self.s_colist_rev:
                    fig0.add_trace(go.Scatter(x=self.dfs["sec"], y=self.dfs[i].astype('float16'), mode='lines', name=i,
                                             line=dict(width=1)), secondary_y=False)
                for j in self.r_colist_rev:
                    fig0.add_trace(go.Scatter(x=self.dfs["sec"], y=self.dfs[j].astype('int8'), mode='lines', name=j,
                                             line=dict(width=1)), secondary_y=True)
                for l in self.lim_C_rev:
                    fig0.add_trace(go.Scatter(x=self.dfs['sec'], y=self.dfs[l], mode='lines', name='Cu Spec',
                                             line=dict(color='#9D303D')))
                for m in self.lim_STS_rev:
                    fig0.add_trace(go.Scatter(x=self.dfs['sec'], y=self.dfs[m], mode='lines', name='STS Spec',
                                             line=dict(color='#9D303D')))

                fig0.update_layout(title_text = "{} {} {} 20{}-{}-{} {}".format(self.model, self.condition, self.sample, self.year, self.month, self.day, self.place),
                    title_font_color = "green",
                    title_font_size = 20,
                    font_family = "LG스마트체 Regular")
                fig0.update_xaxes(title_text="Second(s)", range=[0, len(self.dfs)*2+50], showgrid=False,tickmode = 'linear', tick0 = 0, dtick = 200)
                fig0.update_yaxes(title_text="Stress(MPa)",
                                 range=[-1, self.dfs[self.s_colist].astype(float).max().max() + 5], showgrid=False,
                                 secondary_y=False)

                fig0.update_traces(hovertemplate='%{y:.1f} MPa')

                self.browser.setHtml(fig0.to_html(include_plotlyjs='cdn'))

            elif self.mv0_odu.empty and self.tdiff != 0:
                try:
                    QMessageBox.about(self, 'Warning', '시간차가 없습니다. 0을 입력후 [Show]를 누르세요.')
                except ValueError:
                    pass

            elif self.tdiff != 0:
                try:

                    self.mv0_odu['adj_sec'] = self.mv0_odu['mv_sec'] + self.tdiff

                    # stress_insec = pd.concat([self.dfs['sec'], self.dfs[self.s_colist]], axis=1)

                    tot_df = pd.merge(self.dfs, self.mv0_odu, left_on="sec", right_on="adj_sec")
                    print(tot_df)

                    # s_max = adj_sec.groupby('INV 현재').max()
                    # num = s_max.index.shape[0]
                    #
                    # self.dfs["sec"] = self.dfs.sec.astype(int)
                    # tot_df = pd.merge(self.dfs, self.mv0_odu, left_on="sec", right_on="adj_sec")
                    # tot_df = pd.merge(tot_df, self.mv0_odu)

                    ##############-Save DB-#######################################################

                    path = "./DB"
                    if not os.path.isdir(path):
                        os.mkdir(path)
                    path = os.path.realpath(path)
                    self.DB_name = "./DB/{}_{}_{}.csv".format(self.model, self.condition, self.sample)
                    print(self.DB_name)
                    tot_df.to_csv(self.DB_name, encoding='cp949')

                    ##################################--------------Max List 정리 -----------####################
                    number_place = ((len(self.dfs.columns) - 3) / 2)
                    number_place = int(number_place)
                    pipe_place = self.dfs.columns[1:(number_place + 1)]

                    #----------------------------time comp 운전 계산-----------------------------------
                    read_index = tot_df[tot_df['COMP운전'] > 0].index.tolist()
                    a0 = read_index[0]
                    a1 = read_index[-1]
                    #---------------------------------운전중 Max--------------------------------------
                    tot_df[pipe_place] = tot_df[pipe_place].astype(float)

                    oper_max = []
                    for n in pipe_place:
                        oper_stress = max(tot_df[n][a0:a1])
                        index_o = tot_df[n][a0:a1].index
                        oper_max.append(oper_stress)

                    o_place_index = oper_max.index(max(oper_max))
                    o_place = pipe_place[o_place_index]

                    tot_df_list = []
                    for i in pipe_place:
                        b = []
                        for j in index_o:
                            b.append(tot_df[i][j])
                        tot_df_list.append(b)

                    comp_df_oper = pd.DataFrame(tot_df_list).T
                    comp_df_oper.columns = pipe_place
                    comp_df_oper = comp_df_oper.astype(float)

                    op_max_value = []
                    op_max_idx = []

                    for q in pipe_place:
                        c = max(comp_df_oper[q])
                        op_max_value.append(c)
                        d = comp_df_oper[q].idxmax()
                        op_max_idx.append(d+index_o[0])

                    op_max_comp = []
                    for k in op_max_idx:
                        a = tot_df['INV 현재'][k]
                        op_max_comp.append(a)

                    o_place_index = op_max_value.index(max(op_max_value))
                    o_place = pipe_place[o_place_index]

                    oper_list = {'idx(운전Max)': op_max_idx, '운전Max': op_max_value}
                    oper_max_df = pd.DataFrame(oper_list)

                    oper_max_df = oper_max_df.astype(float)
                    operidxmax = oper_max_df['운전Max'].idxmax()
                    # ---------------------------------시동중 Max--------------------------------------
                    start_max = []
                    for m in pipe_place:
                        start_stress = max(tot_df[m][a0 - 10:a0 + 10])
                        index = tot_df[m][a0 - 10:a0 + 10].index
                        start_max.append(start_stress)

                    s_place_index = start_max.index(max(start_max))
                    s_place = pipe_place[s_place_index]

                    tot_df_a = []
                    for k in index:
                        b1 = tot_df[s_place][k]
                        tot_df_a.append(b1)

                    original_index = tot_df_a.index(max(tot_df_a)) + index[0]
                    # ---------------------------------정지중 Max--------------------------------------
                    finish_max = []

                    for f in pipe_place:
                        finish_stress = max(tot_df[f][a1 - 10:a1 + 10])
                        index_f = tot_df[f][a1 - 10:a1 + 10].index
                        finish_max.append(finish_stress)

                    f_place_index = finish_max.index(max(finish_max))
                    f_place = pipe_place[f_place_index]

                    tot_df_b = []
                    for n in index_f:
                        b2 = tot_df[f_place][n]
                        tot_df_b.append(b2)

                    original_index_f = tot_df_b.index(max(tot_df_b)) + index_f[0]

                    #------------------------------------Max.list table--------------------------------------

                    pipe_place= pipe_place.tolist()
                    start_max = list(map(float, start_max))
                    finish_max = list(map(float, finish_max))
                    oper_max = list(map(float, oper_max))

                    max_list = {'센서위치': pipe_place, '시동Max': start_max, '정지Max': finish_max, '운전Max': oper_max}

                    max_df = pd.DataFrame(max_list)

                    conditionlist = [(max_df['시동Max'] >= 200), (max_df['정지Max'] >= 200), (max_df['운전Max'] >= 90),
                                     (max_df['시동Max'] < 200) & (max_df['시동Max'] >= 40),
                                     (max_df['정지Max'] < 200) & (max_df['정지Max'] >= 40),
                                     (max_df['운전Max'] < 90) & (max_df['운전Max'] >= 15),
                                     (max_df['시동Max'] < 40), (max_df['정지Max'] < 40), (max_df['운전Max'] < 15)]

                    max_df['Comp Hz'] = op_max_comp
                    choicelist = ['NG/NG', 'NG/NG', 'NG/NG', 'NG/OK', 'NG/OK', 'NG/OK', 'OK/OK', 'OK/OK', 'OK/OK']
                    max_df['CU/STS'] = np.select(conditionlist, choicelist, default='Not Specified')

                    max_df2 = max_df
                    max_df2['Comp Hz'] = max_df2['Comp Hz'].astype(int)

                    self.max_df1 = max_df
                    print(self.max_df1)

                    ######------------------------------그래프1---------------------------------------------#

                    fig = make_subplots(specs=[[{"secondary_y": True}]])
                    fig.add_trace(
                        go.Scatter(x=tot_df['sec'], y=tot_df.index, mode='markers', name='Temp', xaxis='x2', visible=False))
                    for k in self.Hz_col_rev:
                        fig.add_trace(go.Scatter(x=tot_df['adj_sec'], y=tot_df[k], mode='lines', name='INV rps',
                                                 line=dict(color="rgb(0,0,0)", width=1)), secondary_y=True)
                    for l in self.lim_C_rev:
                        fig.add_trace(go.Scatter(x=tot_df['sec'], y=tot_df[l], mode='lines', name='Cu Spec',
                                                 line=dict(color='rgb(250,0,0)')))
                    for m in self.lim_STS_rev:
                        fig.add_trace(go.Scatter(x=tot_df['sec'], y=tot_df[m], mode='lines', name='STS Spec',
                                                 line=dict(color='#9D303D')))
                    for i in self.s_colist_rev:
                        fig.add_trace(go.Scatter(x=tot_df["sec"], y=tot_df[i].astype('float16'), mode='lines', name=i,
                                                 line=dict(width=1)), secondary_y=False)
                        fig.update_traces(hovertemplate='%{y:.1f}MPa')

                    for j in self.r_colist_rev:
                        fig.add_trace(go.Scatter(x=tot_df["sec"], y=tot_df[j].astype('int'), mode='lines', name=j,
                                                 line=dict(width=1)), secondary_y=False)
                        fig.update_traces(hovertemplate='%{x:.0f}sec,%{y:.0f}Hz')

                    apply_list = []
                    for i in range(len(fig.data)):
                        apply_list.append(i)

                    time_axis = []
                    for i in range(len(fig.data)):
                        if i == 0:
                            time_axis.append(False)
                        else:
                            time_axis.append(True)

                    Comp_axis = []
                    for i in range(len(fig.data)):
                        if i == 0:
                            Comp_axis.append(False)
                        # elif i == 1:
                        #     Comp_axis.append(False)
                        else:
                            Comp_axis.append(True)

                    button_layer_1_height = 1.1

                    annotations1 = [dict(text="X-axis :", x=0.0, xref="paper", y=1.075, yref="paper", align="left",
                                         showarrow=False, font=dict(color="black", size=15)),
                                    dict(text="Line Colors :", x=0.71, xref="paper", y=1.075, yref="paper",
                                         align="left", showarrow=False, font=dict(color="black", size=15)),
                                    dict(text="Max :", x=0.25, xref="paper", y=1.075, yref="paper",
                                         align="left", showarrow=False, font=dict(color="black", size=15))
                                    ]

                    annotations2 = [dict(text="X-axis : ", x=0.0, xref="paper", y=1.075, yref="paper", align="left",
                                        showarrow=False, font=dict(color="black", size=15)),
                                   dict(text="Line Colors : ", x=0.71, xref="paper", y=1.075, yref="paper",
                                        align="left", showarrow=False, font=dict(color="black", size=15)),
                                   dict(text="Max : ", x=0.25, xref="paper", y=1.075, yref="paper",
                                        align="left", showarrow=False, font=dict(color="black", size=15)),
                                   dict(x=tot_df['sec'][original_index], y=tot_df[s_place][original_index],
                                        text=('시동 Max 위치: {}<br>시동 Max 크기 : {} MPa').format(s_place, tot_df[s_place][
                                            original_index]),
                                        showarrow=True, arrowhead=2, ay=-100, ax=100, arrowsize=1, arrowwidth=2,
                                        arrowcolor="#d7005f"
                                        , font=dict(family='LG스마트체 Light', size=14, color='#d7005f'),
                                        bordercolor="#d7005f", bgcolor="#ffffff", opacity=0.9),

                                   dict(x=tot_df['sec'][original_index_f], y=tot_df[f_place][original_index_f],
                                        text=('정지 Max 위치: {} <br> 정지 Max 크기: {}MPa').format(f_place, tot_df[f_place][
                                            original_index_f]),
                                        showarrow=True, arrowhead=2, ay=30, ax=-100, arrowsize=1, arrowwidth=2,
                                        arrowcolor="#d7005f"
                                        , font=dict(family='LG스마트체 Light', size=14, color='#d7005f'),
                                        bordercolor="#d7005f", bgcolor="#ffffff", opacity=0.9),

                                   dict(x=tot_df['sec'][oper_max_df['idx(운전Max)'][operidxmax]],
                                        y=tot_df[o_place][oper_max_df['idx(운전Max)'][operidxmax]],
                                        text=('운전 Max 위치: {} <br> 운전 Max 크기: {}MPa').format(o_place, tot_df[o_place][
                                            oper_max_df['idx(운전Max)'][operidxmax]]),
                                        showarrow=True, arrowhead=2, ay=-50, ax=-70, arrowsize=1, arrowwidth=2,
                                        arrowcolor="#d7005f"
                                        , font=dict(family='LG스마트체 Light', size=14, color='#d7005f'),
                                        bordercolor="#d7005f", bgcolor="#ffffff", opacity=0.9)
                                   ]


                    layout = go.Layout(
                        plot_bgcolor="#FFFFFF",
                        xaxis=dict(title="Time(s)", linecolor="#BCCCDC", showspikes=True, spikethickness=1, spikedash="dot", spikecolor="#999999", spikemode="across"),
                        xaxis2=dict(visible=False, overlaying='x1', matches='x1',),
                        yaxis=dict(title="Stress(MPa)", linecolor="#BCCCDC", showgrid=True),
                        yaxis2=dict(title="Comp(Hz)", showgrid=False),
                        annotations = annotations1,
                        title_text = ("{} {} {} 20{}-{}-{} {}".format(self.model, self.condition, self.sample, self.year, self.month, self.day, self.place)),
                        title = {'y': 0.97, 'x': 0.2,'xanchor': 'center','yanchor': 'top'},
                        title_font_family = "LG스마트체 Regular", title_font_color = "green", title_font_size = 20, font_family = "LG스마트체 Regular",

                        updatemenus=[
                            dict(buttons=list([
                                dict(label="Label Off", method="relayout",
                                     args=[{'annotations':annotations1}] ),

                                dict(label="Label On", method="relayout",
                                     args=[{'annotations': annotations2}]),
                            ]),
                                direction="down",
                                pad={"r": 10, "t": 8},
                                showactive=True,
                                x=0.5,
                                xanchor="right",
                                y=button_layer_1_height,
                                yanchor="top"),

                            dict(buttons=list([
                                dict(args=[{"marker.color": plotly.colors.DEFAULT_PLOTLY_COLORS}], label="Default",
                                     method="update"),
                                dict(args=[{"marker.color": px.colors.qualitative.Bold}], label="Bold",
                                     method="update"),
                                dict(args=[{"marker.color": px.colors.sequential.Viridis}, apply_list], label="Viridis",
                                     method="update"),
                                dict(args=[{"marker.color": px.colors.sequential.Plasma}, apply_list], label="Plasma",
                                     method="update"),
                                dict(args=[{"marker.color": px.colors.sequential.Cividis}, apply_list], label="Cividis",
                                     method="update"),
                                dict(args=[{"marker.color": px.colors.sequential.Blues}, apply_list], label="Blues",
                                     method="update")]),
                                direction="down", pad={"r": 10, "t": 8}, showactive=True, x=0.72, xanchor="left",
                                y=button_layer_1_height, yanchor="top"),

                            dict(buttons=list([

                                dict(label="Time", method="update",
                                     args=[{"x": [tot_df['sec']], "visible": time_axis},
                                           {"xaxis.tickmode": "auto",
                                            "xaxis.showgrid": False,
                                            "xaxis.showticklabels": True,
                                            "xaxis.title": "time(s)",
                                            "xaxis.nticks": 20,
                                            "xaxis.visible": True,
                                            "xaxis2.visible": False,
                                            "yaxis2.visible": True},
                                           apply_list],
                                     ),
                                dict(label="Comp",
                                     method="update",
                                     args=[{"x": [tot_df['sec']],
                                            "visible": Comp_axis},
                                           {"xaxis.tickmode": "array",
                                            "xaxis.tickvals": tot_df['sec'],
                                            "xaxis.showgrid": False,
                                            "xaxis.showticklabels": False,
                                            "xaxis.ticktext": tot_df['INV 현재'],
                                            "xaxis.title": "",
                                            "yaxis2.visible": False,
                                            "xaxis2.visible": True,
                                            "xaxis2.tickmode": "array",
                                            "xaxis2.tickvals": tot_df['sec'][::100],
                                            "xaxis2.showgrid": False,
                                            "xaxis2.showticklabels": True,
                                            "xaxis2.title": "Comp(Hz)",
                                            "xaxis2.ticktext": tot_df['INV 현재'][::100],
                                            },
                                           apply_list])
                            ]),
                                direction="down",
                                pad={"r": 10, "t": 8},
                                showactive=True,
                                x=0.1,
                                xanchor="left",
                                y=button_layer_1_height,
                                yanchor="top"),
                        ]
                    )

                    fig.update_layout(layout)

                    # #################################------------Comp-응력 그래프----------#################################

                    df_analysis1 = tot_df[pipe_place]
                    df_control = tot_df['INV 현재']
                    df_control = df_control.to_frame()
                    df_analysis1 = pd.concat([df_analysis1, df_control], axis=1)

                    melt_place = pd.melt(df_analysis1, id_vars=['INV 현재'], value_vars=pipe_place)
                    melt_place.columns = ['Comp_Hz','위치','MPa']
                    melt_place.sort_values(by=['Comp_Hz'],axis=0,inplace=True)
                    melt_place1 = melt_place.reset_index(drop=True)
                    melt_place1['MPa'] = melt_place1['MPa'].astype(float)

                    fig1 = px.violin(melt_place1, y='MPa', x ='위치', color='위치', box=True, points='all', hover_data=melt_place.columns,
                                     range_y=[0,(melt_place1['MPa'].max())+3])

                    fig1.update_layout(margin=dict(t=40, b=40, l=10, r=10), template="plotly_white")

                    ##############comp-응력 #########################
                    melt_place1['Comp_Hz']=melt_place1['Comp_Hz'].astype(int)

                    fig2 = px.violin(melt_place1, y='MPa', x ='위치', color='위치', box=True, points='all', hover_data=melt_place.columns,
                                     animation_frame='Comp_Hz', range_y=[0,(melt_place1['MPa'].max())+3])
                    fig2.update_layout(margin=dict(t=40, b=40, l=10, r=10), template="plotly_white")

                    self.browser.setHtml(fig.to_html(include_plotlyjs='cdn'))
                    self.browser1.setHtml(fig1.to_html(include_plotlyjs='cdn'))
                    self.browser2.setHtml(fig2.to_html(include_plotlyjs='cdn'))
                    self.browser3.setHtml(fig.to_html(include_plotlyjs='cdn'))
                    self.browser4.setHtml(fig.to_html(include_plotlyjs='cdn'))
                except ValueError:
                    pass

            else:
                try:
                    QMessageBox.about(self, 'Warning', '3-2.Input Time Diff.\n데이터 시간차를 입력하세요. ex) -20')

                    self.mv0_odu['adj_sec'] = self.mv0_odu['mv_sec'] + self.tdiff

                    s_insec = pd.concat([self.dfs['sec'], self.dfs[self.s_colist]], axis=1)
                    s_insec = s_insec.astype('float')
                    adj_sec = pd.merge(self.mv0_odu, s_insec, left_on="adj_sec", right_on="sec")

                    self.dfs["sec"] = self.dfs.sec.astype(int)
                    tot_df = pd.merge(self.dfs, self.mv0_odu, left_on="sec", right_on="adj_sec")

                    fig0_1 = make_subplots(specs=[[{"secondary_y": True}]])
                    fig0_1.update_layout(margin=dict(t=40, b=40, l=10, r=10), template="plotly_white", )

                    fig0_1.add_trace(
                        go.Scatter(x=tot_df['sec'], y=tot_df.index, mode='markers', name='Temp', xaxis='x2',
                                   visible=False))
                    for k in self.Hz_col_rev:
                        fig0_1.add_trace(go.Scatter(x=tot_df['adj_sec'], y=tot_df[k], mode='lines', name='INV rps',
                                                 line=dict(color="rgb(0,0,0)", width=1)), secondary_y=True)
                    for l in self.lim_C_rev:
                        fig0_1.add_trace(go.Scatter(x=tot_df['sec'], y=tot_df[l], mode='lines', name='Cu Spec',
                                                 line=dict(color='rgb(250,0,0)')))
                    for m in self.lim_STS_rev:
                        fig0_1.add_trace(go.Scatter(x=tot_df['sec'], y=tot_df[m], mode='lines', name='STS Spec',
                                                 line=dict(color='#9D303D')))
                    for i in self.s_colist_rev:
                        fig0_1.add_trace(go.Scatter(x=tot_df["sec"], y=tot_df[i].astype('float16'), mode='lines', name=i,
                                                 line=dict(width=1)), secondary_y=False)
                        fig0_1.update_traces(hovertemplate='%{y:.1f}MPa')

                    for j in self.r_colist_rev:
                        fig0_1.add_trace(go.Scatter(x=tot_df["sec"], y=tot_df[j].astype('int'), mode='lines', name=j,
                                                 line=dict(width=1)), secondary_y=False)
                        fig0_1.update_traces(hovertemplate='%{x:.0f}sec,%{y:.0f}Hz')

                    layout = go.Layout(
                        plot_bgcolor="#FFFFFF",

                        xaxis=dict(
                            title="Time(s)",
                            linecolor="#BCCCDC",
                            showspikes=True,
                            spikethickness=1,
                            spikedash="dot",
                            spikecolor="#999999",
                            spikemode="across",
                        ),
                        xaxis2=dict(
                            visible=False,
                            overlaying='x1',
                            matches='x1',
                        ),
                        yaxis=dict(
                            title="Stress(MPa)",
                            linecolor="#BCCCDC",
                            showgrid=True,
                        ),
                        yaxis2=dict(
                            title="Comp(Hz)",
                            showgrid=False,
                        ),

                        title_text=("{} {} {} 20{}-{}-{} {}".format(self.model, self.condition, self.sample, self.year,
                                                                    self.month, self.day, self.place)),
                        title={
                            'y': 0.97,
                            'x': 0.2,
                            'xanchor': 'center',
                            'yanchor': 'top'},
                        title_font_family="LG스마트체 Regular",
                        title_font_color="green",
                        title_font_size=20,
                        font_family="LG스마트체 Regular")

                    fig0_1.update_layout(layout)

                    self.browser.setHtml(fig0_1.to_html(include_plotlyjs='cdn'))
                except ValueError:
                    pass

    def max_list_show(self):
        try:
            if self.max_df1.empty:
                QMessageBox.about(self, '알림', 'MV데이터와 시간차를 입력해야 Max List를 만들수 있습니다.')
            else:
                row = 0
                self.table_max.setRowCount(len(self.max_df1))
                for id in range(len(self.max_df1)):
                    self.table_max.setItem(row, 0, QtWidgets.QTableWidgetItem(str(self.max_df1['센서위치'][id])))
                    self.table_max.setItem(row, 1, QtWidgets.QTableWidgetItem(str(self.max_df1['CU/STS'][id])))
                    self.table_max.setItem(row, 2, QtWidgets.QTableWidgetItem(str(self.max_df1['Comp Hz'][id])))
                    self.table_max.setItem(row, 3, QtWidgets.QTableWidgetItem(str(self.max_df1['운전Max'][id])))
                    self.table_max.setItem(row, 4, QtWidgets.QTableWidgetItem(str(self.max_df1['시동Max'][id])))
                    self.table_max.setItem(row, 5, QtWidgets.QTableWidgetItem(str(self.max_df1['정지Max'][id])))
                    row = row + 1
        except ValueError:
            pass

    def capture(self):
        self.path = "./CapturedGraphs"
        if not os.path.isdir(self.path):
            os.mkdir(self.path)

        size = self.browser.contentsRect()
        img = QPixmap(size.width(), size.height())
        self.browser.render(img)
        self.cap_name = "./CapturedGraphs/{}_{}_{}.jpg".format(self.model, self.condition, self.sample)
        img.save(self.cap_name)
        QMessageBox.about(self, 'Saved', '그래프가 저장되었습니다. Refresh를 누르세요.')

    def makeppt(self):
        if self.dfs.empty:
            QMessageBox.about(self, 'Warning', '응력 데이터가 없습니다. 응력 데이터를 입력하세요.')
        else:
            self.path = "./report"
            if not os.path.isdir(self.path):
                os.mkdir(self.path)

            model = []
            condition = []
            sample = []
            file_list = self.keywords[0]

            if not file_list:
                file_list = self.file_list
        # for list in file_list:
        #     model.append(list.split("_")[0])
        #     condition.append(list.split("_")[1].split(".")[0])

            for list in file_list:
                model.append(list.split("_")[0])
                condition.append(list.split("_")[1])
                sample.append(list.split("_")[2].split(".")[0])

            prs = Presentation()

        # 제목 페이지
            lyt = prs.slide_layouts[0]
            slide0 = prs.slides.add_slide(lyt)
            title = slide0.shapes.title
            subtitle = slide0.placeholders[1]
            title.text = "{} 응력계측 결과".format(model[0])
            #title.text = "_______ 응력계측 결과"
            subtitle.text = "소음진동 TASK"

        # 두번째 페이지-그림 삽입
            i = 0
            for temp in file_list:
                img_path = './CapturedGraphs/{}'.format(temp)
                blank_slide_layout = prs.slide_layouts[6]
                slide1 = prs.slides.add_slide(blank_slide_layout)
                left = Inches(2)
                top = Inches(1.5)
                width = Inches(7)
                height = Inches(5.5)
                img = slide1.shapes.add_picture(img_path, left, top, width=width, height=height)

                Topline = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.7), Inches(10), Inches(0.005))
                Topline.shadow.inherit = False
            # fill=shape.fill
            # fill.solid()
            # fill.fore_color.rgb=RGBColor(0,0,0)
            # shape.text= "How to Add a Chart"
                line = Topline.line
                line.color.rgb = RGBColor(0, 0, 0)

                text_box1 = slide1.shapes.add_textbox(Inches(0.1), Inches(0.15), Inches(6), Inches(0.5))
                tb1 = text_box1.text_frame
                tb1.text = "{} 응력계측 결과".format(model[i])

                text_box2 = slide1.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(6), Inches(0.5))
                tb2 = text_box2.text_frame
                tb2.text = "■운전조건: {}".format(condition[i])

                text_box3 = slide1.shapes.add_textbox(Inches(3.1), Inches(1.15), Inches(6), Inches(0.5))
                tb3 = text_box3.text_frame
                tb3.text = "■시료: {}".format(sample[i])

                i = i + 1

            prs.save("./report/{}_응력계측결과.pptx".format(model[0]))
            QMessageBox.about(self, 'Completed', 'PPT파일이 생성되었습니다.')

if __name__ == "__main__" :
    #QApplication : 프로그램을 실행시켜주는 클래스
    app = QApplication(sys.argv)
    #WindowClass의 인스턴스 생성
    myWindow = WindowClass()
    #프로그램 화면을 보여주는 코드
    myWindow.show()
    #프로그램을 이벤트루프로 진입시키는(프로그램을 작동시키는) 코드
    app.exec_()
